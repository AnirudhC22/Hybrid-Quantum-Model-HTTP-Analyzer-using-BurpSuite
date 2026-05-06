"""Extract all text from the PPTX file."""
from pptx import Presentation
import sys, io

# Force UTF-8 output
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

pptx_path = r'd:\web_attack_detector\INTELLIGENT WEB REQUEST ATTACK (1) (4).pptx'
out_path = r'd:\web_attack_detector\scratch\pptx_text.txt'
prs = Presentation(pptx_path)

with open(out_path, 'w', encoding='utf-8') as f:
    f.write(f"Total slides: {len(prs.slides)}\n")

    for i, slide in enumerate(prs.slides, 1):
        f.write(f"\n{'='*70}\n")
        f.write(f"SLIDE {i}\n")
        f.write(f"{'='*70}\n")
        
        if slide.slide_layout:
            f.write(f"[Layout: {slide.slide_layout.name}]\n")
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        bold = any(run.font.bold for run in para.runs)
                        if bold:
                            f.write(f"  **{text}**\n")
                        else:
                            f.write(f"  {text}\n")
            
            if shape.has_table:
                table = shape.table
                f.write("\n  [TABLE]\n")
                for row_idx, row in enumerate(table.rows):
                    cells = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
                    f.write(f"  | {' | '.join(cells)} |\n")
                f.write("\n")

print(f"Done. {len(prs.slides)} slides extracted to {out_path}")
